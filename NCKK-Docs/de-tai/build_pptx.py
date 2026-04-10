"""Generate slide-HarnessEval-v2-en.pptx from HarnessEval v2 defense deck.

Run: python build_pptx.py
Output: slide-HarnessEval-v2-en.pptx (19 slides, 16:9)
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from copy import deepcopy

# ---------- Brand ----------
BRAND = RGBColor(0x00, 0x66, 0xCC)
ACCENT = RGBColor(0xCC, 0x33, 0x33)
OK = RGBColor(0x2E, 0x7D, 0x32)
WARN = RGBColor(0xE6, 0x51, 0x00)
MUTED = RGBColor(0x55, 0x55, 0x55)
DARK = RGBColor(0x1A, 0x1A, 0x2E)
BG_CARD = RGBColor(0xFA, 0xFB, 0xFD)
BORDER = RGBColor(0xDC, 0xDF, 0xE6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
HIGHLIGHT_BG = RGBColor(0xFF, 0xF8, 0xE1)
HIGHLIGHT_BORDER = RGBColor(0xFF, 0xA0, 0x00)
RESULT_BG = RGBColor(0xF0, 0xF9, 0xF0)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW = prs.slide_width
SH = prs.slide_height
BLANK = prs.slide_layouts[6]

# ---------- Helpers ----------
def add_slide():
    return prs.slides.add_slide(BLANK)

def add_textbox(slide, left, top, width, height, text, *, size=14, bold=False,
                color=DARK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font='Segoe UI', italic=False):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb

def add_para(textbox, text, *, size=14, bold=False, color=DARK, align=PP_ALIGN.LEFT, font='Segoe UI', italic=False):
    tf = textbox.text_frame
    p = tf.add_paragraph()
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return p

def add_rect(slide, left, top, width, height, fill=BG_CARD, line=BORDER, line_w=0.75):
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    rect.adjustments[0] = 0.08
    rect.fill.solid()
    rect.fill.fore_color.rgb = fill
    rect.line.color.rgb = line
    rect.line.width = Pt(line_w)
    rect.shadow.inherit = False
    rect.text_frame.text = ''
    return rect

def add_header(slide, title, num=None, total=19):
    """Top brand bar + title + slide number."""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(0.7))
    bar.fill.solid(); bar.fill.fore_color.rgb = BRAND
    bar.line.fill.background()
    bar.shadow.inherit = False
    add_textbox(slide, Inches(0.5), Inches(0.1), Inches(11.5), Inches(0.5),
                title, size=22, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    if num is not None:
        add_textbox(slide, Inches(12.0), Inches(0.1), Inches(1.2), Inches(0.5),
                    f"{num} / {total}", size=12, color=WHITE, align=PP_ALIGN.RIGHT,
                    anchor=MSO_ANCHOR.MIDDLE)

def add_footer(slide):
    add_textbox(slide, Inches(0.5), Inches(7.15), Inches(12.0), Inches(0.3),
                "HarnessEval v2  |  FPT School of Business and Technology (FSB)  |  April 2026",
                size=9, color=MUTED)

def add_card(slide, left, top, width, height, title, body, *, title_color=BRAND, border=BORDER):
    rect = add_rect(slide, left, top, width, height, fill=BG_CARD, line=border, line_w=1.0)
    tf = rect.text_frame
    tf.margin_left = Inches(0.15); tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.12); tf.margin_bottom = Inches(0.12)
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = title
    r.font.name = 'Segoe UI'
    r.font.size = Pt(14)
    r.font.bold = True
    r.font.color.rgb = title_color
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.LEFT
    r2 = p2.add_run()
    r2.text = body
    r2.font.name = 'Segoe UI'
    r2.font.size = Pt(11)
    r2.font.color.rgb = DARK
    return rect

def add_kpi(slide, left, top, width, height, num, label):
    rect = add_rect(slide, left, top, width, height, fill=BG_CARD, line=BORDER, line_w=1.0)
    tf = rect.text_frame
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.15); tf.margin_bottom = Inches(0.05)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = num
    r.font.name = 'Segoe UI'
    r.font.size = Pt(36)
    r.font.bold = True
    r.font.color.rgb = BRAND
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = label
    r2.font.name = 'Segoe UI'
    r2.font.size = Pt(10)
    r2.font.color.rgb = MUTED
    return rect

def add_note(slide, left, top, width, height, text, *, fill=RGBColor(0xF0, 0xF7, 0xFF), border=BRAND):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    rect.fill.solid(); rect.fill.fore_color.rgb = fill
    rect.line.color.rgb = border
    rect.line.width = Pt(1.0)
    rect.shadow.inherit = False
    # Left accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.06), height)
    bar.fill.solid(); bar.fill.fore_color.rgb = border
    bar.line.fill.background()
    bar.shadow.inherit = False
    tf = rect.text_frame
    tf.margin_left = Inches(0.18); tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.1); tf.margin_bottom = Inches(0.1)
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = text
    r.font.name = 'Segoe UI'
    r.font.size = Pt(11)
    r.font.color.rgb = DARK
    return rect

def add_table(slide, left, top, width, height, data, *, header_fill=RGBColor(0xF0, 0xF2, 0xF5),
              first_col_bold=False, font_size=10):
    rows = len(data)
    cols = len(data[0])
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    tbl = table_shape.table
    for r_idx, row in enumerate(data):
        for c_idx, val in enumerate(row):
            cell = tbl.cell(r_idx, c_idx)
            cell.margin_left = Inches(0.08)
            cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.04)
            cell.margin_bottom = Inches(0.04)
            tf = cell.text_frame
            tf.word_wrap = True
            tf.clear()
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = str(val)
            run.font.name = 'Segoe UI'
            run.font.size = Pt(font_size)
            if r_idx == 0:
                run.font.bold = True
                run.font.color.rgb = DARK
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_fill
            else:
                if first_col_bold and c_idx == 0:
                    run.font.bold = True
                run.font.color.rgb = DARK
                if r_idx % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(0xF8, 0xF9, 0xFC)
                else:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = WHITE
    return table_shape

# ===========================================================
# SLIDE 1 — TITLE
# ===========================================================
s = add_slide()
# Background brand stripe
stripe = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.4), SW, Inches(0.04))
stripe.fill.solid(); stripe.fill.fore_color.rgb = BRAND
stripe.line.fill.background()
stripe.shadow.inherit = False

add_textbox(s, Inches(0.5), Inches(1.0), Inches(12.3), Inches(1.4),
            "HarnessEval", size=64, bold=True, color=BRAND, align=PP_ALIGN.CENTER)
add_textbox(s, Inches(0.5), Inches(2.6), Inches(12.3), Inches(1.0),
            "Evaluating Coding Agent Scaffolds Through",
            size=22, color=DARK, align=PP_ALIGN.CENTER)
add_textbox(s, Inches(0.5), Inches(3.05), Inches(12.3), Inches(0.6),
            "Ablation-Driven Multi-Dimensional Analysis",
            size=22, color=DARK, align=PP_ALIGN.CENTER)
add_textbox(s, Inches(0.5), Inches(3.85), Inches(12.3), Inches(0.5),
            "A unified framework for decomposing harness quality from base LLM capability",
            size=13, color=MUTED, align=PP_ALIGN.CENTER)

add_textbox(s, Inches(0.5), Inches(5.0), Inches(12.3), Inches(0.5),
            "Methods of Learning & Scientific Research",
            size=14, color=MUTED, align=PP_ALIGN.CENTER)
add_textbox(s, Inches(0.5), Inches(5.35), Inches(12.3), Inches(0.5),
            "FPT School of Business and Technology (FSB)",
            size=16, bold=True, color=DARK, align=PP_ALIGN.CENTER)
add_textbox(s, Inches(0.5), Inches(5.75), Inches(12.3), Inches(0.5),
            "April 2026",
            size=12, color=MUTED, align=PP_ALIGN.CENTER)

# v2 badge
badge = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.9), Inches(6.4), Inches(1.5), Inches(0.45))
badge.adjustments[0] = 0.5
badge.fill.solid(); badge.fill.fore_color.rgb = RGBColor(0xE3, 0xF2, 0xFD)
badge.line.color.rgb = BRAND; badge.line.width = Pt(1.0)
badge.shadow.inherit = False
btf = badge.text_frame
btf.margin_left = Inches(0); btf.margin_right = Inches(0)
btf.margin_top = Inches(0); btf.margin_bottom = Inches(0)
btf.vertical_anchor = MSO_ANCHOR.MIDDLE
bp = btf.paragraphs[0]; bp.alignment = PP_ALIGN.CENTER
br = bp.add_run(); br.text = "v2 — revised"
br.font.name = 'Segoe UI'; br.font.size = Pt(11); br.font.bold = True
br.font.color.rgb = BRAND

# ===========================================================
# SLIDE 2 — TEAM
# ===========================================================
s = add_slide()
add_header(s, "Research Team", 2)
add_textbox(s, Inches(0.5), Inches(1.0), Inches(12.3), Inches(0.5),
            "Supervisor", size=18, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
add_textbox(s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(0.5),
            "[Supervisor name — please provide]",
            size=18, bold=True, color=DARK, align=PP_ALIGN.CENTER)
add_textbox(s, Inches(0.5), Inches(1.95), Inches(12.3), Inches(0.4),
            "FPT School of Business and Technology (FSB)",
            size=12, color=MUTED, align=PP_ALIGN.CENTER)

add_textbox(s, Inches(0.5), Inches(2.7), Inches(12.3), Inches(0.5),
            "Team Members", size=18, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

team_data = [
    ["#", "Full Name", "Student ID"],
    ["1", "Lê Lâm Vĩnh", "25MS23328"],
    ["2", "Hoàng Phi Hải", "25MS23323"],
    ["3", "Đỗ Hoàng Tỷ Phú", "25MS23327"],
    ["4", "Trịnh Hữu Tuấn", "25MS23315"],
    ["5", "Trần Quang Tuấn", "25MS23331"],
]
add_table(s, Inches(3.5), Inches(3.3), Inches(6.3), Inches(3.3), team_data, font_size=14)
add_footer(s)

# ===========================================================
# SLIDE 3 — OUTLINE
# ===========================================================
s = add_slide()
add_header(s, "Outline", 3)
left_items = [
    "1. Problem — the evaluation paradox",
    "2. Research gap — three reasons",
    "3. Aim & 5 objectives",
    "4. Hypotheses H1–H4",
    "5. Taxonomy: 3 dimensions, 7 metrics",
    "6. Modular Harness architecture",
    "7. Methodology — 4 phases",
]
right_items = [
    "8. Ablation design 27×150",
    "9. Two-way ANOVA",
    "10. Threats & Risks",
    "11. Timeline & budget",
    "12. Implementation status (167 tests)",
    "13. Preliminary empirical results",
    "14. Contributions & publishability",
    "15. Q&A",
]
tb = add_textbox(s, Inches(0.8), Inches(1.2), Inches(6.0), Inches(5.5),
                 left_items[0], size=18, color=DARK)
for item in left_items[1:]:
    add_para(tb, item, size=18, color=DARK)
tb2 = add_textbox(s, Inches(7.0), Inches(1.2), Inches(5.8), Inches(5.5),
                  right_items[0], size=18, color=DARK)
for item in right_items[1:]:
    add_para(tb2, item, size=18, color=DARK)
add_footer(s)

# ===========================================================
# SLIDE 4 — PROBLEM
# ===========================================================
s = add_slide()
add_header(s, "1. Problem: The Evaluation Paradox", 4)
add_textbox(s, Inches(0.5), Inches(1.0), Inches(12.3), Inches(0.5),
            "The performance of a coding agent depends on two factors:",
            size=16, color=DARK)
add_card(s, Inches(0.7), Inches(1.6), Inches(5.7), Inches(1.4),
         "A. Base LLM",
         "Claude, GPT, DeepSeek, Qwen…\nMany benchmarks already evaluate this.")
add_card(s, Inches(6.9), Inches(1.6), Inches(5.7), Inches(1.4),
         "B. Harness / Scaffold",
         "Tool system, context manager, orchestration…\n\"The operating system of the agent\" [A1]")

add_note(s, Inches(0.7), Inches(3.2), Inches(11.9), Inches(1.5),
         "Evidence: Lou et al. (2026) [A2] — a small model with a good harness can beat a large model "
         "with a weak harness, reducing invalid actions by 100%.  Bui et al. (2026) [A1] — Adaptive "
         "Context Compaction cuts token usage by 54%.",
         fill=HIGHLIGHT_BG, border=HIGHLIGHT_BORDER)

add_note(s, Inches(0.7), Inches(4.9), Inches(11.9), Inches(1.5),
         "The paradox: when we report \"Agent A scored 70%,\" we do NOT know which harness component "
         "contributed to that result.",
         fill=RGBColor(0xFF, 0xF8, 0xF0), border=WARN)
add_footer(s)

# ===========================================================
# SLIDE 5 — RESEARCH GAP
# ===========================================================
s = add_slide()
add_header(s, "2. Research Gap: Three Genuine Reasons", 5)
add_textbox(s, Inches(0.5), Inches(1.0), Inches(12.3), Inches(0.5),
            "Analysis of 49 papers (2023–2026): no work systematically evaluates the harness itself [S5].",
            size=14, color=DARK, italic=False)

card_w = Inches(4.0); card_h = Inches(3.2); top = Inches(1.7); gap = Inches(0.2)
add_card(s, Inches(0.6), top, card_w, card_h,
         "1. Role mismatch",
         "Princeton (SWE-Agent), Anthropic (Claude Code) are developers — incentivized to ship "
         "new harnesses, not to compare systems as evaluators.\n\nAnalogy: car makers vs Euro NCAP / IIHS")
add_card(s, Inches(4.7), top, card_w, card_h,
         "2. Field moving too fast",
         "2023–2026: from a handful of harnesses to dozens. Nobody has paused to meta-evaluate "
         "while still racing to build.")
add_card(s, Inches(8.8), top, card_w, card_h,
         "3. Only indirect evidence",
         "Lou et al. compared harnesses on game tasks [A2]. Robeyns et al. self-improved "
         "a harness but did NOT measure each component [B10].")

add_note(s, Inches(0.6), Inches(5.2), Inches(12.2), Inches(1.5),
         "Research gap: No framework exists to (a) measure the quality of individual harness components, "
         "(b) decompose harness vs base-LLM contribution, (c) compare harnesses on unified quality dimensions.",
         fill=HIGHLIGHT_BG, border=HIGHLIGHT_BORDER)
add_footer(s)

# ===========================================================
# SLIDE 6 — GAP TABLE
# ===========================================================
s = add_slide()
add_header(s, "2.1. What Exists vs What Is Missing", 6)
gap_data = [
    ["Aspect", "What exists", "What is missing"],
    ["Measuring agent output", "SWE-bench, SWE-Compass, AgentBoard", "—"],
    ["Tool accuracy in isolation", "ToolLLM, AnyTool, Gorilla", "Tool use INSIDE a coding harness"],
    ["Memory in isolation", "Mem0, A-MEM (on chatbots)", "Memory INSIDE a coding harness"],
    ["Harness design description", "OpenDev, AutoHarness", "Measuring harness quality"],
    ["Cross-backend comparison", "—", "Completely unaddressed"],
    ["Separating harness vs LLM", "—", "Completely unaddressed"],
]
add_table(s, Inches(0.5), Inches(1.2), Inches(12.3), Inches(5.5), gap_data, font_size=14)
add_footer(s)

# ===========================================================
# SLIDE 7 — AIM + OBJECTIVES
# ===========================================================
s = add_slide()
add_header(s, "3. Aim + 5 Objectives", 7)
add_note(s, Inches(0.5), Inches(1.0), Inches(12.3), Inches(1.4),
         "Aim: Build the FIRST UNIFIED EVALUATION FRAMEWORK for coding-agent harnesses, enabling "
         "decomposition and quantification of each harness component's contribution to agent "
         "performance, INDEPENDENT of the base LLM.",
         fill=HIGHLIGHT_BG, border=HIGHLIGHT_BORDER)

obj_data = [
    ["Obj.", "Description", "Success criterion"],
    ["O1", "Taxonomy: 3 dimensions, 7 metrics, validated by 5 experts", "Fleiss κ ≥ 0.6"],
    ["O2", "Fork & refactor SWE-Agent into a modular harness", "± 3% of original SWE-Agent"],
    ["O3", "Ablation: 150 tasks × 27 conditions (3×3×3) + pilot study", "7,050 evaluations"],
    ["O4", "Two-way ANOVA decomposing harness vs LLM variance", "p < 0.05 (Bonferroni)"],
    ["O5", "Release open-source toolkit on GitHub", "Docs + CI/CD"],
]
add_table(s, Inches(0.5), Inches(2.7), Inches(12.3), Inches(3.4), obj_data, first_col_bold=True, font_size=13)
add_note(s, Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.7),
         "v2 changes — Reduced 5 dimensions → 3, 80 → 27 conditions, added pilot study, added two-way ANOVA.",
         fill=RGBColor(0xE3, 0xF2, 0xFD), border=BRAND)
add_footer(s)

# ===========================================================
# SLIDE 8 — HYPOTHESES
# ===========================================================
s = add_slide()
add_header(s, "4. Research Hypotheses", 8)
ch = Inches(2.3); cw = Inches(6.0)
add_card(s, Inches(0.5), Inches(1.1), cw, ch, "H1",
         "Tool system has the LARGEST effect size (Cohen's d > 0.5) among the three harness "
         "components on resolve rate.")
add_card(s, Inches(6.8), Inches(1.1), cw, ch, "H2",
         "Context management has a MEDIUM effect size (Cohen's d 0.3 – 0.5).")
add_card(s, Inches(0.5), Inches(3.6), cw, ch, "H3",
         "Harness configuration explains AT LEAST 20% of resolve-rate variance (η² ≥ 0.20), "
         "independent of the LLM backend.")
add_card(s, Inches(6.8), Inches(3.6), cw, ch, "H4",
         "Variance across LLM backends DECREASES when harness quality is higher (statistically "
         "significant interaction effect).")

add_note(s, Inches(0.5), Inches(6.1), Inches(12.3), Inches(0.9),
         "Neutral framing — no specific numbers predicted. Only testable directions → avoid confirmation bias. "
         "Supporting OR refuting a hypothesis is equally valuable.",
         fill=RESULT_BG, border=OK)
add_footer(s)

# ===========================================================
# SLIDE 9 — TAXONOMY
# ===========================================================
s = add_slide()
add_header(s, "5. Taxonomy: 3 Dimensions — 7 Metrics", 9)
tax_data = [
    ["Dim.", "Metric", "Definition", "Measurement"],
    ["D1\nTool", "M1.1 Correct Selection Rate", "% calls picking from acceptable tool set", "2 annotators, Fleiss κ"],
    ["D1", "M1.2 Redundant Call Rate", "% calls whose output unused in next 3 turns", "Auto log analysis"],
    ["D1", "M1.3 Utilization Breadth", "unique tools used / total available", "Auto log analysis"],
    ["D2\nContext", "M2.1 Info Retention", "BERTScore: full vs compacted output", "30 human-preference pairs"],
    ["D2", "M2.2 Effective Token Ratio", "% tokens relevant to current task", "LLM classifier + 200 labels"],
    ["D3\nBackend", "M3.1 Cross-Backend StdDev", "std(resolve rate) across 3 LLMs", "Direct"],
    ["D3", "M3.2 Min/Max Ratio", "min(RR) / max(RR)", "Higher = more portable"],
]
add_table(s, Inches(0.4), Inches(1.1), Inches(12.5), Inches(5.0), tax_data, font_size=11)
add_note(s, Inches(0.4), Inches(6.2), Inches(12.5), Inches(0.7),
         "v2 — M1.1 changed \"single ground truth\" → acceptable-tool set (committee Q#4). "
         "M2.2 changed vague \"token waste\" → LLM classifier + human validation (Q#9).",
         fill=RGBColor(0xE3, 0xF2, 0xFD), border=BRAND)
add_footer(s)

# ===========================================================
# SLIDE 10 — MODULAR HARNESS
# ===========================================================
s = add_slide()
add_header(s, "6. Modular Harness — 3 Swappable Layers", 10)

def add_layer(slide, top, height, color, title, sub):
    rect = add_rect(slide, Inches(1.5), top, Inches(10.3), height,
                    fill=BG_CARD, line=color, line_w=2.0)
    tf = rect.text_frame
    tf.margin_left = Inches(0.2); tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.1); tf.margin_bottom = Inches(0.1)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = title
    r.font.name = 'Segoe UI'; r.font.size = Pt(16); r.font.bold = True
    r.font.color.rgb = color
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.LEFT
    r2 = p2.add_run(); r2.text = sub
    r2.font.name = 'Segoe UI'; r2.font.size = Pt(12); r2.font.color.rgb = MUTED

add_layer(s, Inches(1.1), Inches(1.0), ACCENT,
          "LAYER 1: Tool Dispatch — ToolProvider",
          "Full (12 tools)   |   Medium (8)   |   Minimal (5)")
add_layer(s, Inches(2.25), Inches(1.0), WARN,
          "LAYER 2: Context Management — ContextStrategy",
          "Full History   |   Sliding Window (50K)   |   Summary / ACC (2K)")
add_layer(s, Inches(3.4), Inches(1.0), OK,
          "LAYER 3: LLM Backend — LLMBackend",
          "Claude Sonnet 4   |   GPT-4o   |   DeepSeek-V3")

# Muted future layers
add_textbox(s, Inches(1.5), Inches(4.6), Inches(10.3), Inches(0.4),
            "L4 Safety & Permission — future work", size=11, color=MUTED, italic=True)
add_textbox(s, Inches(1.5), Inches(5.0), Inches(10.3), Inches(0.4),
            "L5 Orchestration — kept unchanged from SWE-Agent", size=11, color=MUTED, italic=True)

add_note(s, Inches(0.5), Inches(5.7), Inches(12.3), Inches(1.2),
         "Techniques: Strategy Pattern + Dependency Injection + Factory Pattern. Each layer is "
         "swapped via a YAML config. Forked from SWE-Agent (MIT, ~15k LOC).")
add_footer(s)

# ===========================================================
# SLIDE 11 — METHOD
# ===========================================================
s = add_slide()
add_header(s, "7. Methodology — Four Experimental Phases", 11)
phase_w = Inches(2.95); phase_h = Inches(3.0); top = Inches(1.3)
add_card(s, Inches(0.4), top, phase_w, phase_h,
         "P1 — Taxonomy",
         "5 experts validate metrics.\nFleiss κ ≥ 0.6.\n\nWeeks 1–6")
add_card(s, Inches(3.5), top, phase_w, phase_h,
         "P1.5 — Pilot (v2)",
         "5 cond × 20 tasks × 2 runs = 200 evals.\nCost ~$80.\n\nWeeks 7–8")
add_card(s, Inches(6.6), top, phase_w, phase_h,
         "P2 — Modular Fork",
         "SWE-Agent refactor.\nVerify ±3% on 50 tasks.\n\nWeeks 9–11")
add_card(s, Inches(9.7), top, phase_w, phase_h,
         "P3 — Full Ablation",
         "27 cond × 150 tasks = 7,050 evals.\n\nWeeks 12–18")

add_note(s, Inches(0.4), Inches(4.7), Inches(12.5), Inches(2.0),
         "Why a pilot? Validate that metrics actually discriminate between conditions, estimate cost "
         "accurately, debug the pipeline, and compute preliminary effect sizes for power analysis. "
         "Prevents BLIND COMMITMENT to the expensive full experiment.")
add_footer(s)

# ===========================================================
# SLIDE 12 — ABLATION DESIGN
# ===========================================================
s = add_slide()
add_header(s, "8. Ablation Design: 27 = 3 × 3 × 3", 12)

kpi_top = Inches(1.1); kpi_w = Inches(2.95); kpi_h = Inches(1.4)
add_kpi(s, Inches(0.4), kpi_top, kpi_w, kpi_h, "27", "CONDITIONS")
add_kpi(s, Inches(3.5), kpi_top, kpi_w, kpi_h, "150", "TASKS")
add_kpi(s, Inches(6.6), kpi_top, kpi_w, kpi_h, "7,050", "EVALUATIONS")
add_kpi(s, Inches(9.7), kpi_top, kpi_w, kpi_h, "~$2.5k", "BUDGET")

factor_data = [
    ["Factor", "Level 1", "Level 2", "Level 3"],
    ["Tool Config", "Full (12 tools)", "Medium (8)", "Minimal (5)"],
    ["Context Strategy", "Full History", "Sliding Window 50K", "Summary / ACC 2K"],
    ["LLM Backend", "Claude Sonnet 4", "GPT-4o", "DeepSeek-V3"],
]
add_table(s, Inches(0.4), Inches(2.9), Inches(12.5), Inches(2.0), factor_data, first_col_bold=True, font_size=13)

add_note(s, Inches(0.4), Inches(5.2), Inches(12.5), Inches(1.5),
         "Dataset: 150 tasks from SWE-bench Verified, stratified random sampling by difficulty.\n"
         "Runs: 10 critical conditions × 3 runs + 17 conditions × 1 run.\n"
         "Power analysis: α = 0.05, power ≥ 0.80.")
add_footer(s)

# ===========================================================
# SLIDE 13 — ANOVA
# ===========================================================
s = add_slide()
add_header(s, "9. Two-Way ANOVA — Decomposing Harness vs LLM", 13)

add_note(s, Inches(0.5), Inches(1.0), Inches(12.3), Inches(1.2),
         "\"Independent of the LLM\" does NOT mean harness and LLM are uncorrelated. It means we can "
         "DECOMPOSE the variance due to harness vs the variance due to LLM.",
         fill=HIGHLIGHT_BG, border=HIGHLIGHT_BORDER)

# Code-style block
code_box = slide_code = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(2.5), Inches(12.3), Inches(2.8))
code_box.fill.solid(); code_box.fill.fore_color.rgb = RGBColor(0xF5, 0xF7, 0xFA)
code_box.line.color.rgb = BORDER; code_box.line.width = Pt(0.75)
code_box.shadow.inherit = False
ctf = code_box.text_frame
ctf.margin_left = Inches(0.2); ctf.margin_right = Inches(0.2)
ctf.margin_top = Inches(0.15); ctf.margin_bottom = Inches(0.15)
ctf.word_wrap = True
lines = [
    "Total Variance = Var(Harness) + Var(LLM) + Var(Harness × LLM) + Error",
    "",
    "Factor A: Harness config   (9 levels: 3 tools × 3 contexts)",
    "Factor B: LLM backend      (3 levels: Claude, GPT, DeepSeek)",
    "DV:        Resolve rate",
    "",
    "If Var(Harness) significant   →  harness has an INDEPENDENT effect",
    "If Var(H × LLM) significant   →  there is an INTERACTION",
    "η² (eta-squared)              →  % of variance explained per factor",
]
for i, line in enumerate(lines):
    if i == 0:
        p = ctf.paragraphs[0]
    else:
        p = ctf.add_paragraph()
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = line
    r.font.name = 'Consolas'
    r.font.size = Pt(13)
    r.font.color.rgb = DARK

add_note(s, Inches(0.5), Inches(5.5), Inches(12.3), Inches(1.4),
         "Main theoretical contribution: formalizing how to SEPARATE harness quality from LLM quality. "
         "Post-hoc: Tukey HSD, Cohen's d, Bonferroni correction for 7 metrics.",
         fill=RESULT_BG, border=OK)
add_footer(s)

# ===========================================================
# SLIDE 14 — THREATS & RISKS
# ===========================================================
s = add_slide()
add_header(s, "10. Threats to Validity & Risk Analysis", 14)

add_textbox(s, Inches(0.5), Inches(1.0), Inches(6.0), Inches(0.4),
            "Threats", size=16, bold=True, color=ACCENT)
threats = [
    ["Type", "Mitigation"],
    ["Internal: tool-removal confound", "Reduce 12→8→5; log per-task tool needs"],
    ["Internal: LLM non-determinism", "Temperature = 0; 3 runs critical conds"],
    ["External: Python only", "Caution when generalizing"],
    ["Construct: M1.1 annotator-dependent", "2 annotators, Fleiss κ"],
]
add_table(s, Inches(0.5), Inches(1.5), Inches(6.0), Inches(3.5), threats, font_size=11)

add_textbox(s, Inches(6.8), Inches(1.0), Inches(6.0), Inches(0.4),
            "Top Risks", size=16, bold=True, color=ACCENT)
risks = [
    ["Risk", "P", "Mitigation"],
    ["API cost overrun", "High", "Pilot estimates; DeepSeek cheap"],
    ["Docker env instability", "Med", "Pre-filter, retry ×2"],
    ["Metrics not discriminative", "Med", "Detected early via pilot"],
    ["DeepSeek API unstable", "Low", "Backup: Qwen2.5-72B"],
]
add_table(s, Inches(6.8), Inches(1.5), Inches(6.0), Inches(3.5), risks, font_size=11)

add_note(s, Inches(0.5), Inches(5.4), Inches(12.3), Inches(1.3),
         "v2 — V1 had no Threats or Risk Analysis. The committee explicitly requested this section.",
         fill=RGBColor(0xFF, 0xF8, 0xF0), border=WARN)
add_footer(s)

# ===========================================================
# SLIDE 15 — TIMELINE
# ===========================================================
s = add_slide()
add_header(s, "11. Timeline (24 weeks) & Budget", 15)
timeline = [
    ["Phase", "Weeks", "Content", "Deliverable"],
    ["P1", "1–6", "Taxonomy + expert validation", "Taxonomy v1.0, κ report"],
    ["P1.5", "7–8", "Pilot study (200 evals)", "Metrics validation, cost"],
    ["P2", "9–11", "Fork SWE-Agent + modular refactor", "Modular harness v1.0"],
    ["P3", "12–18", "Full ablation (7,050 evals)", "Raw data, ANOVA results"],
    ["P4", "19–22", "Analysis + guidelines + thesis", "Draft thesis"],
    ["P5", "23–24", "Review, revisions, toolkit release", "Final + GitHub release"],
]
add_table(s, Inches(0.5), Inches(1.1), Inches(12.3), Inches(3.5), timeline, first_col_bold=True, font_size=12)

add_card(s, Inches(0.5), Inches(4.9), Inches(6.0), Inches(2.0),
         "Estimated budget",
         "Pilot:               ~$80\nFull experiment:     ~$2,500\nContingency:         ~$500\n\n"
         "TOTAL:               $2,500 – $3,100")
add_card(s, Inches(6.8), Inches(4.9), Inches(6.0), Inches(2.0),
         "v1 vs v2 comparison",
         "v1: $800–1,200 (unrealistic)\nv2: $2,500–3,100 (incl. retry + pilot)\n\n"
         "Evaluations: 16,000 → 7,050 (-56%)")
add_footer(s)

# ===========================================================
# SLIDE 16 — IMPLEMENTATION STATUS
# ===========================================================
s = add_slide()
add_header(s, "12. Implementation Status — ALREADY BUILT", 16)

kpi_top = Inches(1.1); kpi_w = Inches(2.95); kpi_h = Inches(1.4)
add_kpi(s, Inches(0.4), kpi_top, kpi_w, kpi_h, "167", "TESTS PASSING")
add_kpi(s, Inches(3.5), kpi_top, kpi_w, kpi_h, "9", "PHASES DONE")
add_kpi(s, Inches(6.6), kpi_top, kpi_w, kpi_h, "27", "CONDITIONS")
add_kpi(s, Inches(9.7), kpi_top, kpi_w, kpi_h, "360", "TRAJECTORIES")

add_textbox(s, Inches(0.5), Inches(2.8), Inches(6.0), Inches(0.4),
            "Toolkit  harness_eval/", size=14, bold=True, color=ACCENT)
toolkit_items = [
    "• configs/  — 27 conditions (3×3×3)",
    "• metrics/  — 7 metrics M1.1–M3.2",
    "• parsers/trajectory.py — SWE-Agent log parser",
    "• pipeline/analysis.py  — ANOVA, Cohen's d, GLMM",
    "• pipeline/runner.py    — Dry-run + SWE-Agent + Ollama",
    "• harness/  — 3 ABCs + 9 providers + factory",
    "• cli.py    — info / pilot / run / analyze / convert",
]
tb = add_textbox(s, Inches(0.5), Inches(3.25), Inches(6.0), Inches(3.6),
                 toolkit_items[0], size=11, color=DARK)
for it in toolkit_items[1:]:
    add_para(tb, it, size=11, color=DARK)

add_textbox(s, Inches(6.8), Inches(2.8), Inches(6.0), Inches(0.4),
            "Streamlit Dashboard", size=14, bold=True, color=ACCENT)
dash_items = [
    "• 5 tabs: Config / Run Monitor / Logs / Compare / ANOVA",
    "• Run modes: ollama / dry-run / real",
    "• Parallel ThreadPoolExecutor (4 workers)",
    "• 150 SWE-bench tasks embedded (difficulty-aware)",
    "• Ollama: qwen2.5:7b, deepseek-r1:1.5b",
    "• One-click setup: start_harness_eval.sh",
]
tb2 = add_textbox(s, Inches(6.8), Inches(3.25), Inches(6.0), Inches(3.6),
                  dash_items[0], size=11, color=DARK)
for it in dash_items[1:]:
    add_para(tb2, it, size=11, color=DARK)
add_footer(s)

# ===========================================================
# SLIDE 17 — PRELIMINARY RESULTS
# ===========================================================
s = add_slide()
add_header(s, "13. Preliminary Results — Empirical ANOVA", 17)
add_textbox(s, Inches(0.5), Inches(1.0), Inches(12.3), Inches(0.4),
            "Run on current trajectories/: 360 observations, 18 conditions (2 backends: claude + sonnet)",
            size=12, color=MUTED, italic=True)

anova = [
    ["Source", "η²", "F", "p-value", "Sig?"],
    ["Tool", "0.0539", "10.264", "< 0.001", "***"],
    ["Context", "0.0248", "4.723", "0.0095", "**"],
    ["Backend", "0.0025", "0.959", "0.3282", "—"],
    ["Tool × Context", "0.0058", "0.550", "0.6988", "—"],
    ["Tool × Backend", "0.0069", "1.314", "0.2701", "—"],
    ["Context × Backend", "0.0069", "1.314", "0.2701", "—"],
    ["Error", "0.8979", "—", "—", "—"],
]
tbl = add_table(s, Inches(0.5), Inches(1.5), Inches(7.5), Inches(3.6), anova, first_col_bold=True, font_size=11)
# Highlight Tool and Context rows green
for r_idx in (1, 2):
    for c_idx in range(5):
        cell = tbl.table.cell(r_idx, c_idx)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE9)

# Result box
add_note(s, Inches(8.2), Inches(1.5), Inches(4.7), Inches(3.6),
         "Preliminary observations:\n\n"
         "• Tool has the largest effect (η²=0.054 > Context=0.025) — H1 direction supported.\n\n"
         "• Context has a medium, statistically significant effect — H2 direction supported.\n\n"
         "• Backend has only 2 levels (haiku empty); H4 needs the full design.",
         fill=RESULT_BG, border=OK)

add_textbox(s, Inches(0.5), Inches(5.4), Inches(12.3), Inches(0.5),
            "Note: current data is realistic dry-run, not yet real SWE-bench runs. "
            "Empirical results follow after Phase P3.",
            size=11, color=MUTED, italic=True)
add_footer(s)

# ===========================================================
# SLIDE 18 — CONTRIBUTIONS + PUBLISHABILITY
# ===========================================================
s = add_slide()
add_header(s, "14. Contributions & Publishability", 18)
cw = Inches(4.0); ch = Inches(2.0); top = Inches(1.1)
add_card(s, Inches(0.4), top, cw, ch,
         "Theoretical",
         "Formalizing how to separate harness quality from LLM quality via two-way ANOVA.")
add_card(s, Inches(4.6), top, cw, ch,
         "Empirical",
         "First quantitative decomposition of each harness component (Tool / Context / Backend).")
add_card(s, Inches(8.8), top, cw, ch,
         "Toolkit",
         "Reusable evaluation framework — open-source on GitHub, docs, CI/CD.")

add_textbox(s, Inches(0.5), Inches(3.4), Inches(12.3), Inches(0.4),
            "Practical impact", size=14, bold=True, color=ACCENT)
practical = [
    "• Developers — know which harness component to prioritize for improvement",
    "• Researchers — standard framework to compare new harnesses",
    "• Industry — data-driven guidelines for harness design",
]
tb = add_textbox(s, Inches(0.6), Inches(3.85), Inches(12.0), Inches(1.2),
                 practical[0], size=12, color=DARK)
for it in practical[1:]:
    add_para(tb, it, size=12, color=DARK)

add_textbox(s, Inches(0.5), Inches(5.2), Inches(12.3), Inches(0.4),
            "Publishability", size=14, bold=True, color=ACCENT)
pub = [
    ["Venue", "Type", "Fit"],
    ["ICSE 2027 NIER", "Conference (New Ideas)", "Strong fit"],
    ["MSR 2027", "Conference", "Fit"],
    ["LLM Agents Workshop (NeurIPS/ICML)", "Workshop", "Fit"],
    ["EMNLP Demo Track", "Demo", "Fit"],
]
add_table(s, Inches(0.5), Inches(5.65), Inches(12.3), Inches(1.4), pub, font_size=11)
add_footer(s)

# ===========================================================
# SLIDE 19 — Q&A
# ===========================================================
s = add_slide()
# Soft background
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
bg.fill.solid(); bg.fill.fore_color.rgb = RGBColor(0xF5, 0xF7, 0xFA)
bg.line.fill.background()
bg.shadow.inherit = False

add_textbox(s, Inches(0.5), Inches(1.8), Inches(12.3), Inches(1.8),
            "Q & A", size=120, bold=True, color=BRAND, align=PP_ALIGN.CENTER)
add_textbox(s, Inches(0.5), Inches(4.0), Inches(12.3), Inches(0.6),
            "Thank you for your attention",
            size=22, color=DARK, align=PP_ALIGN.CENTER)
add_textbox(s, Inches(0.5), Inches(4.5), Inches(12.3), Inches(0.6),
            "and thoughtful questions",
            size=22, color=DARK, align=PP_ALIGN.CENTER)

add_textbox(s, Inches(0.5), Inches(5.6), Inches(12.3), Inches(0.4),
            "Resources", size=12, bold=True, color=MUTED, align=PP_ALIGN.CENTER)
resources = [
    "NCKK-Docs/de-tai/de-cuong/DE-CUONG-HarnessEval-v2.md",
    "NCKK-Docs/de-tai/harness-eval/  — toolkit + 167 tests",
    "plans/Plan-Coding-HarnessEval.md  — session log",
]
tb = add_textbox(s, Inches(0.5), Inches(5.95), Inches(12.3), Inches(1.2),
                 resources[0], size=11, color=MUTED, align=PP_ALIGN.CENTER, font='Consolas')
for r in resources[1:]:
    add_para(tb, r, size=11, color=MUTED, align=PP_ALIGN.CENTER, font='Consolas')

# ---------- Save ----------
out = "slide-HarnessEval-v2-en.pptx"
prs.save(out)
print(f"Saved: {out} ({len(prs.slides)} slides)")
